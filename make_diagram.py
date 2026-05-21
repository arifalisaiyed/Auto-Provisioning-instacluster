"""Generates architecture diagram as PNG and PPTX."""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ───────────────────────────────────────────────────────────────────
BG   = "#0f172a"
C_TITLE = "#e2e8f0"
C_SUB   = "#94a3b8"
C_LABEL = "#cbd5e1"

# fill colours
F_PROD = "#0c2340"
F_KFKA = "#2e1065"
F_CONS = "#052e16"
F_GROQ = "#451a03"
F_CASS = "#0c2340"
F_DASH = "#0f172a"
F_TF   = "#1e293b"

# border / accent colours
B_PROD = "#38bdf8"
B_KFKA = "#a78bfa"
B_CONS = "#4ade80"
B_GROQ = "#fbbf24"
B_CASS = "#67e8f9"
B_DASH = "#818cf8"
B_TF   = "#94a3b8"
B_ARR  = "#475569"

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


# ═══════════════════════════════════════════════════════════════════════════════
#  PNG via matplotlib
# ═══════════════════════════════════════════════════════════════════════════════

FW, FH = 20, 12
fig, ax = plt.subplots(figsize=(FW, FH))
fig.patch.set_facecolor(BG)
ax.set_facecolor(BG)
ax.set_xlim(0, FW);  ax.set_ylim(0, FH)
ax.axis("off")


def mbox(ax, x, y, w, h, title, sub, fill, border, tsz=9.5, ssz=7.8):
    ax.add_patch(FancyBboxPatch((x, y), w, h,
        boxstyle="round,pad=0.04", lw=1.8,
        edgecolor=border, facecolor=fill, zorder=3))
    cy = y + h/2
    offset = 0.13 if sub else 0
    ax.text(x+w/2, cy+offset, title,
            ha="center", va="center", fontsize=tsz,
            fontweight="bold", color=C_TITLE, zorder=4)
    if sub:
        ax.text(x+w/2, cy-0.18, sub,
                ha="center", va="center", fontsize=ssz,
                color=C_SUB, zorder=4, linespacing=1.45)


def grect(ax, x, y, w, h, label, border):
    ax.add_patch(FancyBboxPatch((x, y), w, h,
        boxstyle="round,pad=0.06", lw=1.1,
        linestyle="--", edgecolor=border,
        facecolor="none", zorder=1, alpha=0.5))
    ax.text(x+0.12, y+h-0.07, label,
            fontsize=7.5, color=border, zorder=2, va="top", fontweight="bold")


def arr(ax, x1, y1, x2, y2, label="", color=B_ARR, rad=0.0):
    ax.annotate("", xy=(x2,y2), xytext=(x1,y1),
        arrowprops=dict(arrowstyle="-|>", color=color, lw=1.5,
                        connectionstyle=f"arc3,rad={rad}"), zorder=2)
    if label:
        mx, my = (x1+x2)/2, (y1+y2)/2
        ax.text(mx+0.06, my, label, fontsize=6.5, color=C_SUB, zorder=5, va="center",
                bbox=dict(boxstyle="round,pad=0.18", fc=BG, ec="none", alpha=0.9))


# ── Title ──────────────────────────────────────────────────────────────────────
ax.text(FW/2, FH-0.28, "Infra Monitor  —  Component Architecture",
        ha="center", va="top", fontsize=18, fontweight="bold", color=C_TITLE)
ax.text(FW/2, FH-0.72, "AI-powered infrastructure monitoring  ·  NetApp Instaclustr  ·  Groq LLM",
        ha="center", va="top", fontsize=10.5, color=C_SUB)

# ── TERRAFORM group ─────────────────────────────────────────────────────────
grect(ax, 0.3, 0.3, 4.2, 2.0, "TERRAFORM  ·  instaclustr/instaclustr v2", B_TF)
mbox(ax, 0.55, 0.55, 1.75, 1.4, "Kafka Cluster",
     "KFK-DEV-t4g.small-30\n3 nodes  ·  KRaft", F_TF, B_TF, 8.5, 7.5)
mbox(ax, 2.5,  0.55, 1.75, 1.4, "Cassandra Cluster",
     "CAS-DEV-t4g.medium-30\n3 nodes  ·  RF=3", F_TF, B_TF, 8.5, 7.5)

# ── KAFKA (Instaclustr) ────────────────────────────────────────────────────
grect(ax, 5.0, 0.3, 5.2, 2.0, "APACHE KAFKA 3.9.1  ·  NetApp Instaclustr  (AWS us-east-1)", B_KFKA)
mbox(ax, 5.25, 0.55, 4.7, 1.4, "Topic: server-stats",
     "3 brokers  ·  SASL_PLAINTEXT  ·  SCRAM-SHA-256 auth\nauto.offset.reset = latest  ·  3 partitions  ·  RF=3",
     F_KFKA, B_KFKA, 10, 8)

# ── CASSANDRA (Instaclustr) ───────────────────────────────────────────────
grect(ax, 10.7, 0.3, 5.2, 2.0, "APACHE CASSANDRA 4.1.9  ·  NetApp Instaclustr  (AWS us-east-1)", B_CASS)
mbox(ax, 10.95, 1.55, 4.7, 0.6, "server_stats",
     "time-series  ·  TTL 1 day  ·  PK(server_id, collected_at)", F_CASS, B_CASS, 8.5, 7.5)
mbox(ax, 10.95, 0.97, 4.7, 0.55, "server_latest",
     "one row per server  ·  PK(server_id)", F_CASS, B_CASS, 8.5, 7.5)
mbox(ax, 10.95, 0.40, 4.7, 0.55, "remediation_actions",
     "TTL 7 days  ·  PK(server_id, actioned_at)", F_CASS, B_CASS, 8.5, 7.5)

# ── PRODUCER ──────────────────────────────────────────────────────────────
grect(ax, 0.3, 2.7, 4.2, 5.5, "PRODUCER CONTAINER  ·  Python", B_PROD)
mbox(ax, 0.55, 7.2, 3.7, 0.75, "APScheduler",
     "BlockingScheduler  ·  fires every 5 minutes", F_PROD, B_PROD)
mbox(ax, 0.55, 6.2, 3.7, 0.85, "stats_generator.py",
     "5 servers (web-01/02, db-01/02, cache-01)\n20% anomaly rate  ·  cpu · mem · disk · net", F_PROD, B_PROD)
mbox(ax, 0.55, 5.1, 3.7, 0.9, "kafka_publisher.py",
     "confluent-kafka  ·  JSON payload\nserver_id as message key  ·  flush after batch", F_PROD, B_PROD)
mbox(ax, 0.55, 3.9, 3.7, 1.0, "Simulated Servers",
     "web-01  ·  web-02  ·  db-01  ·  db-02  ·  cache-01\nbaseline ranges + anomaly injection per run", F_PROD, B_PROD, 8.5, 7.5)

arr(ax, 2.4, 7.2,  2.4, 7.05, color=B_PROD)
arr(ax, 2.4, 6.2,  2.4, 6.0,  color=B_PROD)
arr(ax, 2.4, 5.1,  2.4, 4.9,  color=B_PROD)

# ── CONSUMER ──────────────────────────────────────────────────────────────
grect(ax, 5.0, 2.7, 5.2, 7.0, "CONSUMER CONTAINER  ·  Python", B_CONS)
mbox(ax, 5.25, 8.6, 4.7, 0.8, "confluent-kafka Consumer",
     "group: infra-monitor-consumer  ·  poll loop  ·  auto-commit", F_CONS, B_CONS)
mbox(ax, 5.25, 7.5, 4.7, 0.85, "thresholds.py",
     "CPU > 85%  ·  Memory > 90%  ·  Disk > 95%  ·  Net In > 40 MB", F_CONS, B_CONS)
mbox(ax, 5.25, 6.4, 4.7, 0.85, "llm_diagnose.py",
     "LangChain  ·  ChatGroq  ·  StrOutputParser\nprompt: stats + triggered thresholds", F_CONS, B_CONS)
mbox(ax, 5.25, 5.3, 4.7, 0.85, "actions.py",
     "restart_server  ·  clear_temp_files\nscale_up_resources  ·  run_hello_world_script", F_CONS, B_CONS)
mbox(ax, 5.25, 4.2, 4.7, 0.85, "cassandra_writer.py",
     "cassandra-driver  ·  upsert server_stats + server_latest\ninsert remediation_actions on breach only", F_CONS, B_CONS)
mbox(ax, 5.25, 3.0, 4.7, 0.95, "python-dotenv  ·  logging",
     "env-driven config  ·  structured log output", F_CONS, B_CONS, 8.5, 7.5)

arr(ax, 7.6, 8.6, 7.6, 8.35, color=B_CONS)
arr(ax, 7.6, 7.5, 7.6, 7.25, color=B_CONS)
arr(ax, 7.6, 6.4, 7.6, 6.15, color=B_CONS)
arr(ax, 7.6, 5.3, 7.6, 5.05, color=B_CONS)
arr(ax, 7.6, 4.2, 7.6, 3.95, color=B_CONS)

# ── GROQ ──────────────────────────────────────────────────────────────────
grect(ax, 10.7, 7.5, 5.2, 2.2, "GROQ CLOUD API", B_GROQ)
mbox(ax, 10.95, 7.75, 4.7, 1.7, "llama-3.3-70b-versatile",
     "Receives: server metrics + triggered threshold names\nResponds: diagnosis paragraph + ACTION keyword\nHTTPS REST  ·  ~1 second latency  ·  LangChain chain",
     F_GROQ, B_GROQ, 10, 8)

# ── DASHBOARD ─────────────────────────────────────────────────────────────
grect(ax, 10.7, 2.7, 5.2, 4.6, "DASHBOARD CONTAINER  ·  Python", B_DASH)
mbox(ax, 10.95, 6.1, 4.7, 0.8, "FastAPI + Uvicorn",
     "GET /  ·  GET /api/status  ·  GET /stream/kafka\nport 8000", F_DASH, B_DASH)
mbox(ax, 10.95, 5.05, 4.7, 0.8, "cassandra_reader.py",
     "SELECT server_latest  ·  SELECT remediation_actions\nmerge + sort across partitions", F_DASH, B_DASH)
mbox(ax, 10.95, 4.0, 4.7, 0.8, "kafka_stream.py  ·  SSE endpoint",
     "confluent-kafka  ·  group: dashboard-kafka-stream\nyields text/event-stream to browser", F_DASH, B_DASH)
mbox(ax, 10.95, 3.1, 2.2, 0.7, "Overview Tab",
     "Status + Actions\n30s auto-refresh", F_DASH, B_DASH, 8, 7)
mbox(ax, 13.4, 3.1, 2.2, 0.7, "Live Stream Tab",
     "EventSource · SSE\npause + clear", F_DASH, B_DASH, 8, 7)

arr(ax, 13.35, 6.1,  13.35, 5.85, color=B_DASH)
arr(ax, 13.35, 5.05, 13.35, 4.80, color=B_DASH)
arr(ax, 13.35, 4.0,  13.35, 3.80, color=B_DASH)

# ── CROSS-COMPONENT ARROWS ────────────────────────────────────────────────

# Producer -> Kafka (publish)
arr(ax, 4.25, 5.55, 5.25, 1.25, label="publish JSON", color=B_PROD, rad=-0.15)

# Kafka -> Consumer (poll)
arr(ax, 7.6, 2.3, 7.6, 8.6, label="poll messages", color=B_KFKA, rad=0.0)

# Consumer llm_diagnose -> Groq (request)
arr(ax, 10.2, 6.83, 10.7, 8.85, label="HTTPS request", color=B_GROQ, rad=-0.2)
# Groq -> Consumer (response)
arr(ax, 10.7, 8.5, 10.2, 7.0, label="diagnosis + action", color=B_GROQ, rad=0.2)

# Consumer cassandra_writer -> Cassandra (write)
arr(ax, 10.2, 4.62, 10.95, 1.8, label="write", color=B_CASS, rad=-0.1)
arr(ax, 10.2, 4.62, 10.95, 1.24, color=B_CASS, rad=-0.05)
arr(ax, 10.2, 4.62, 10.95, 0.68, color=B_CASS, rad=0.0)

# Dashboard cassandra_reader -> Cassandra (SELECT)
arr(ax, 15.65, 5.05, 13.5, 1.62, label="SELECT", color=B_CASS, rad=0.2)

# Dashboard kafka_stream -> Kafka (poll SSE)
arr(ax, 10.95, 4.4, 10.2, 1.25, label="poll (SSE)", color=B_KFKA, rad=0.15)

# Terraform -> Kafka provision
arr(ax, 2.3, 2.3, 6.5, 2.3, label="provision", color=B_TF)
arr(ax, 3.4, 2.3, 13.35, 2.3, color=B_TF)

# ── LEGEND ────────────────────────────────────────────────────────────────
legend = [
    (B_PROD, F_PROD, "Producer  ·  APScheduler + confluent-kafka"),
    (B_KFKA, F_KFKA, "Kafka (NetApp Instaclustr)  ·  SCRAM-SHA-256"),
    (B_CONS, F_CONS, "Consumer  ·  LangChain + cassandra-driver"),
    (B_GROQ, F_GROQ, "Groq Cloud API  ·  llama-3.3-70b-versatile"),
    (B_CASS, F_CASS, "Cassandra (NetApp Instaclustr)  ·  cassandra-driver"),
    (B_DASH, F_DASH, "Dashboard  ·  FastAPI + Jinja2 + SSE"),
    (B_TF,   F_TF,   "Terraform  ·  instaclustr/instaclustr v2"),
]
cols = 4
for i, (border, fill, label) in enumerate(legend):
    col, row = i % cols, i // cols
    lx = 0.4 + col * 4.9
    ly = 0.02 + row * 0.22 - (0 if row == 0 else 0)
    # tiny swatch
    ax.add_patch(FancyBboxPatch((lx, ly), 0.3, 0.17,
        boxstyle="round,pad=0.02", lw=1.2,
        edgecolor=border, facecolor=fill, zorder=5))
    ax.text(lx+0.38, ly+0.085, label,
            fontsize=7, color=C_LABEL, va="center", zorder=5)

plt.tight_layout(pad=0.2)
fig.savefig("architecture.png", dpi=150, bbox_inches="tight",
            facecolor=BG, edgecolor="none")
print("Saved architecture.png")
plt.close()


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = rgb(BG)


def ptxt(slide, text, left, top, w, h, size=10, bold=False,
         color="#e2e8f0", align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def pbox(slide, title, sub, left, top, w, h,
         fill_h, border_h, tsz=9, ssz=7.5):
    sh = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill_h)
    sh.line.color.rgb = rgb(border_h)
    sh.line.width = Pt(1.6)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(5)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(tsz)
    r.font.bold = True
    r.font.color.rgb = rgb("#e2e8f0")

    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(ssz)
        r2.font.color.rgb = rgb("#94a3b8")


def pconn(slide, x1, y1, x2, y2, color_h=B_ARR):
    c = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(color_h)
    c.line.width = Pt(1.5)


# Title
ptxt(slide, "Infra Monitor  —  Component Architecture",
     0, 0.05, 13.33, 0.38, size=22, bold=True)
ptxt(slide, "AI-powered infrastructure monitoring  ·  NetApp Instaclustr  ·  Groq LLM",
     0, 0.42, 13.33, 0.28, size=10.5, color="#94a3b8")

# ── row 1: Terraform + Kafka + Cassandra ──────────────────────────────────
# Terraform
ptxt(slide, "TERRAFORM  ·  instaclustr/instaclustr v2",
     0.15, 0.78, 3.5, 0.22, size=7, color=B_TF, align=PP_ALIGN.LEFT)
pbox(slide, "Kafka Cluster",     "KFK-DEV-t4g.small-30 · 3 nodes · KRaft",
     0.15, 0.98, 1.6, 0.78, F_TF, B_TF, 8.5, 7.5)
pbox(slide, "Cassandra Cluster", "CAS-DEV-t4g.medium-30 · 3 nodes · RF=3",
     1.9,  0.98, 1.7, 0.78, F_TF, B_TF, 8.5, 7.5)

# Kafka
ptxt(slide, "APACHE KAFKA 3.9.1  ·  NetApp Instaclustr",
     3.75, 0.78, 5.0, 0.22, size=7, color=B_KFKA, align=PP_ALIGN.LEFT)
pbox(slide, "Topic: server-stats",
     "3 brokers · SASL_PLAINTEXT · SCRAM-SHA-256 auth\nauto.offset.reset = latest · 3 partitions · RF=3",
     3.75, 0.98, 5.0, 0.78, F_KFKA, B_KFKA, 10, 8)

# Cassandra
ptxt(slide, "APACHE CASSANDRA 4.1.9  ·  NetApp Instaclustr",
     8.95, 0.78, 4.2, 0.22, size=7, color=B_CASS, align=PP_ALIGN.LEFT)
pbox(slide, "server_stats",          "time-series · TTL 1d · PK(server_id, collected_at)",
     8.95, 0.98, 4.2, 0.46, F_CASS, B_CASS, 8, 7)
pbox(slide, "server_latest",         "one row per server · PK(server_id)",
     8.95, 1.47, 4.2, 0.46, F_CASS, B_CASS, 8, 7)
pbox(slide, "remediation_actions",   "TTL 7d · PK(server_id, actioned_at)",
     8.95, 1.96, 4.2, 0.46, F_CASS, B_CASS, 8, 7)

# ── row 2–4: Producer + Consumer + Groq + Dashboard ──────────────────────
# Producer
ptxt(slide, "PRODUCER CONTAINER  ·  Python",
     0.15, 2.55, 3.5, 0.22, size=7, color=B_PROD, align=PP_ALIGN.LEFT)
pbox(slide, "APScheduler",        "fires every 5 minutes",
     0.15, 2.75, 3.5, 0.55, F_PROD, B_PROD)
pbox(slide, "stats_generator.py", "5 servers · 20% anomaly rate · cpu · mem · disk · net",
     0.15, 3.34, 3.5, 0.65, F_PROD, B_PROD)
pbox(slide, "kafka_publisher.py", "confluent-kafka · JSON · server_id key · flush after batch",
     0.15, 4.03, 3.5, 0.65, F_PROD, B_PROD)
pbox(slide, "Simulated Servers",
     "web-01 · web-02 · db-01 · db-02 · cache-01",
     0.15, 4.72, 3.5, 0.65, F_PROD, B_PROD)
pconn(slide, 1.9, 2.75, 1.9, 3.34, B_PROD)
pconn(slide, 1.9, 3.34, 1.9, 4.03, B_PROD)
pconn(slide, 1.9, 4.03, 1.9, 4.72, B_PROD)

# Consumer
ptxt(slide, "CONSUMER CONTAINER  ·  Python",
     3.75, 2.55, 5.0, 0.22, size=7, color=B_CONS, align=PP_ALIGN.LEFT)
pbox(slide, "confluent-kafka consumer",
     "group: infra-monitor-consumer · poll loop · auto-commit",
     3.75, 2.75, 5.0, 0.6, F_CONS, B_CONS)
pbox(slide, "thresholds.py",
     "CPU > 85%  ·  Memory > 90%  ·  Disk > 95%  ·  Net In > 40 MB",
     3.75, 3.39, 5.0, 0.6, F_CONS, B_CONS)
pbox(slide, "llm_diagnose.py",
     "LangChain · ChatGroq · StrOutputParser",
     3.75, 4.03, 5.0, 0.6, F_CONS, B_CONS)
pbox(slide, "actions.py",
     "restart_server · clear_temp_files · scale_up_resources · run_hello_world_script",
     3.75, 4.67, 5.0, 0.6, F_CONS, B_CONS)
pbox(slide, "cassandra_writer.py",
     "cassandra-driver · upsert server_stats + server_latest · insert remediation_actions on breach",
     3.75, 5.31, 5.0, 0.65, F_CONS, B_CONS)
pconn(slide, 6.25, 2.75, 6.25, 3.39, B_CONS)
pconn(slide, 6.25, 3.39, 6.25, 4.03, B_CONS)
pconn(slide, 6.25, 4.03, 6.25, 4.67, B_CONS)
pconn(slide, 6.25, 4.67, 6.25, 5.31, B_CONS)

# Groq
ptxt(slide, "GROQ CLOUD API",
     3.75, 6.1, 5.0, 0.22, size=7, color=B_GROQ, align=PP_ALIGN.LEFT)
pbox(slide, "llama-3.3-70b-versatile",
     "Input: server metrics + triggered thresholds\nOutput: diagnosis paragraph + ACTION keyword · HTTPS REST",
     3.75, 6.3, 5.0, 0.88, F_GROQ, B_GROQ, 10, 8)

# Dashboard
ptxt(slide, "DASHBOARD CONTAINER  ·  Python",
     8.95, 2.55, 4.2, 0.22, size=7, color=B_DASH, align=PP_ALIGN.LEFT)
pbox(slide, "FastAPI + Uvicorn",
     "GET /  ·  GET /api/status  ·  GET /stream/kafka  ·  port 8000",
     8.95, 2.75, 4.2, 0.6, F_DASH, B_DASH)
pbox(slide, "cassandra_reader.py",
     "SELECT server_latest · SELECT remediation_actions",
     8.95, 3.39, 4.2, 0.6, F_DASH, B_DASH)
pbox(slide, "kafka_stream.py  (SSE)",
     "confluent-kafka · group: dashboard-kafka-stream · text/event-stream",
     8.95, 4.03, 4.2, 0.6, F_DASH, B_DASH)
pbox(slide, "Overview Tab",
     "Server Status + Action History\nauto-refresh 30s",
     8.95, 4.72, 2.0, 0.75, F_DASH, B_DASH, 8, 7)
pbox(slide, "Live Stream Tab",
     "EventSource · SSE\npause + clear",
     11.1, 4.72, 2.0, 0.75, F_DASH, B_DASH, 8, 7)
pbox(slide, "Bootstrap 5  ·  dark theme",  "",
     8.95, 5.51, 4.2, 0.45, F_DASH, B_DASH, 8, 7)
pconn(slide, 11.05, 2.75, 11.05, 3.39, B_DASH)
pconn(slide, 11.05, 3.39, 11.05, 4.03, B_DASH)
pconn(slide, 11.05, 4.03, 11.05, 4.72, B_DASH)

# ── cross-component connectors ────────────────────────────────────────────
# Producer -> Kafka
pconn(slide, 3.65, 4.36, 3.75, 1.37, B_PROD)
# Kafka -> Consumer
pconn(slide, 6.25, 1.76, 6.25, 2.75, B_KFKA)
# Consumer llm -> Groq (request + response)
pconn(slide, 7.0, 4.33, 6.25, 6.3,  B_GROQ)
pconn(slide, 6.25, 6.3, 7.0, 4.93, B_GROQ)
# Consumer cassandra_writer -> Cassandra
pconn(slide, 8.75, 5.63, 8.95, 1.21, B_CASS)
pconn(slide, 8.75, 5.63, 8.95, 1.70, B_CASS)
pconn(slide, 8.75, 5.63, 8.95, 2.19, B_CASS)
# Dashboard cassandra_reader -> Cassandra
pconn(slide, 8.95, 3.69, 13.15, 1.70, B_CASS)
# Dashboard kafka_stream -> Kafka topic
pconn(slide, 8.95, 4.33, 8.75, 1.37, B_KFKA)
# Terraform -> Kafka + Cassandra (provision)
pconn(slide, 3.65, 1.37, 3.75, 1.37, B_TF)
pconn(slide, 3.65, 1.37, 8.95, 1.37, B_TF)

prs.save("architecture.pptx")
print("Saved architecture.pptx")
